"""
replace_ppt_figures.py
Reemplaza las imágenes de los slides 14 y 15 del PPT de defensa
con las versiones corregidas (etiquetas en español).
"""

import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree
import copy
import os

PPT_PATH = r"C:\Users\sonia\Desktop\2_TFM_MELANOMA\PRESENTACION\pptx final\TFM_Salmeron_Fructuoso_Sonia.pptx"
FIG_DIR  = r"C:\Users\sonia\Desktop\2_TFM_MELANOMA\TFM_memoria\memoria_latex\figures"

# Slide indices (0-based) y sus nuevas figuras
REPLACEMENTS = {
    13: os.path.join(FIG_DIR, "km_stage_group_plot.png"),        # slide 14
    14: os.path.join(FIG_DIR, "cox_multivariate_forest_plot.png") # slide 15
}

def replace_largest_image(slide, new_img_path):
    """Localiza la imagen más grande del slide y la reemplaza en sitio."""
    best = None
    best_area = 0

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            area = shape.width * shape.height
            if area > best_area:
                best_area = area
                best = shape

    if best is None:
        print(f"  AVISO: No se encontro ninguna imagen en el slide.")
        return False

    # Guardar posición y tamaño
    left   = best.left
    top    = best.top
    width  = best.width
    height = best.height

    # Añadir la nueva imagen con misma posición/tamaño
    pic = slide.shapes.add_picture(new_img_path, left, top, width, height)

    # Eliminar la imagen antigua del árbol XML
    sp_tree = slide.shapes._spTree
    sp_tree.remove(best._element)

    # Mover el nuevo elemento a la posición original en el árbol
    # (para mantener el z-order aproximado)
    sp_tree.remove(pic._element)
    sp_tree.append(pic._element)

    print(f"  OK: Imagen reemplazada: {os.path.basename(new_img_path)}")
    print(f"    Posicion: left={left}, top={top}, width={width}, height={height}")
    return True


def main():
    print(f"Abriendo: {PPT_PATH}")
    prs = Presentation(PPT_PATH)

    for slide_idx, new_img_path in REPLACEMENTS.items():
        slide = prs.slides[slide_idx]
        print(f"\nSlide {slide_idx + 1} ({slide.shapes.title.text if slide.shapes.title else 'sin título'}):")
        if not os.path.exists(new_img_path):
            print(f"  ERROR: Figura no encontrada: {new_img_path}")
            sys.exit(1)
        replace_largest_image(slide, new_img_path)

    prs.save(PPT_PATH)
    print(f"\nPPT guardado: {PPT_PATH}")


if __name__ == "__main__":
    main()
