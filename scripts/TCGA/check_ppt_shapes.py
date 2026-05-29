"""Check shapes in slides 14 and 15 of the PPT."""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPT_PATH = r"C:\Users\sonia\Desktop\2_TFM_MELANOMA\PRESENTACION\pptx final\TFM_Salmeron_Fructuoso_Sonia.pptx"

prs = Presentation(PPT_PATH)

for idx in [13, 14]:
    slide = prs.slides[idx]
    print(f"\n=== Slide {idx+1} ===")
    for i, shape in enumerate(slide.shapes):
        print(f"  Shape {i}: type={shape.shape_type}, name={shape.name}, "
              f"left={shape.left}, top={shape.top}, "
              f"width={shape.width}, height={shape.height}")
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print(f"    -> PICTURE: {shape.image.content_type}, size={len(shape.image.blob)} bytes")
